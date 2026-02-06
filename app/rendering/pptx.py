"""PPTX rendering — generate PowerPoint from structured JSON data.

All python-pptx layout logic lives here.
"""

from __future__ import annotations

import contextlib
import io
import math
import re
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Cm, Pt

from app.core.log import safe_print


def create_pptx(
    json_data: dict[str, Any],
    template_pptx_bytes: bytes | None = None,
) -> io.BytesIO:
    """Generate a ``.pptx`` from *json_data* and return as BytesIO."""
    if not isinstance(json_data, dict):
        raise ValueError("Dữ liệu slide không hợp lệ (expected dict).")

    # ── Load / create presentation ──────────────────────────────────────
    if template_pptx_bytes:
        try:
            prs = Presentation(io.BytesIO(template_pptx_bytes))
        except Exception as exc:
            safe_print(f"Template PPTX lỗi, dùng mặc định: {exc}")
            prs = Presentation()
        # Clear existing slides from template
        xml_slides = prs.slides._sldIdLst
        for s in list(xml_slides):
            prs.part.drop_rel(s.rId)
            xml_slides.remove(s)
        safe_print(f"Template loaded and cleared. Remaining slides: {len(prs.slides)}")
    else:
        prs = Presentation()

    # ── Layout helpers ──────────────────────────────────────────────────
    def _get_layout(preferred_index: int, needs_body: bool = False):
        if preferred_index < len(prs.slide_layouts):
            layout = prs.slide_layouts[preferred_index]
            if not needs_body or len(layout.placeholders) > 1:
                return layout
        for layout in prs.slide_layouts:
            if needs_body and len(layout.placeholders) > 1:
                return layout
        return prs.slide_layouts[0]

    # ── Title slide ─────────────────────────────────────────────────────
    title_layout = _get_layout(0)
    slide = prs.slides.add_slide(title_layout)
    if slide.shapes.title:
        slide.shapes.title.text = json_data.get("title", "Bài thuyết trình AI").upper()
    if len(slide.placeholders) > 1:
        with contextlib.suppress(IndexError, KeyError):
            slide.placeholders[1].text = "Được tạo bởi SlideGenius"

    # ── Content slides ──────────────────────────────────────────────────
    content_layout = _get_layout(1, needs_body=True)

    for slide_data in json_data.get("slides", []):
        slide = prs.slides.add_slide(content_layout)
        title_height = Cm(0)

        # --- Title ---
        if slide.shapes.title:
            raw_title = slide_data.get("title", "")
            clean_title = re.sub(r"^Slide\s+\d+[:.]?\s*", "", raw_title, flags=re.IGNORECASE)
            slide.shapes.title.text = clean_title

            font_size_pt = 36
            if (
                slide.shapes.title.text_frame
                and slide.shapes.title.text_frame.paragraphs
                and slide.shapes.title.text_frame.paragraphs[0].font.size
            ):
                font_size_pt = slide.shapes.title.text_frame.paragraphs[0].font.size.pt

            slide.shapes.title.left = Cm(1.0)
            slide.shapes.title.width = prs.slide_width - Cm(2.0)
            slide.shapes.title.top = Cm(0.5)

            avg_char_width = font_size_pt * 0.55
            chars_per_line = slide.shapes.title.width.pt / avg_char_width
            estimated_lines = max(1, math.ceil(len(clean_title) / chars_per_line))
            title_height = Pt(estimated_lines * font_size_pt * 1.1)
            slide.shapes.title.height = title_height

            if slide.shapes.title.text_frame:
                tf = slide.shapes.title.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                tf.vertical_anchor = MSO_ANCHOR.TOP
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                    p.font.size = Pt(font_size_pt)
                    p.font.bold = True
                    p.alignment = PP_ALIGN.LEFT
                    tf.margin_top = 0
                    tf.margin_bottom = 0

        # --- Body ---
        body_shape = _find_body_shape(slide)
        if body_shape and hasattr(body_shape, "text_frame"):
            _fill_body(prs, slide, body_shape, slide_data, title_height)

        # --- Notes ---
        notes_text = slide_data.get("notes", "")
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    # ── Save ────────────────────────────────────────────────────────────
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ── Internal helpers ─────────────────────────────────────────────────────


def _find_body_shape(slide):
    """Locate the best body placeholder on *slide*."""
    candidates = []
    for shape in slide.placeholders:
        pht = shape.placeholder_format.type
        if pht in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE):
            continue
        if pht in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
            candidates.append((0, shape))
        elif shape.placeholder_format.idx == 1:
            candidates.append((1, shape))
        else:
            candidates.append((2, shape))
    if candidates:
        candidates.sort(key=lambda x: x[0])
        return candidates[0][1]
    # Fallback: anything that isn't the title
    if len(slide.placeholders) > 1:
        for shape in slide.placeholders:
            if slide.shapes.title is None or shape.element is not slide.shapes.title.element:
                return shape
    return None


def _fill_body(prs, slide, body_shape, slide_data, title_height):
    """Populate the body placeholder with styled content."""
    margin_side = Cm(1.5)
    gap = Cm(0.4)
    margin_top = Cm(0.5) + title_height + gap
    margin_bottom = Cm(1.0)

    body_shape.left = margin_side
    body_shape.top = margin_top
    body_shape.width = prs.slide_width - (margin_side * 2)
    available_height = prs.slide_height - margin_top - margin_bottom
    body_shape.height = available_height

    tf = body_shape.text_frame
    tf.word_wrap = True

    MAX_FONT_SIZE = 24
    MIN_FONT_SIZE = 10
    BASE_CAPACITY_AT_24PT = 300

    total_text_len = sum(len(str(c)) for c in slide_data.get("content", []))
    if total_text_len <= BASE_CAPACITY_AT_24PT:
        font_size_pt = MAX_FONT_SIZE
    else:
        ratio = BASE_CAPACITY_AT_24PT / total_text_len
        font_size_pt = max(min(MAX_FONT_SIZE * math.sqrt(ratio), MAX_FONT_SIZE), MIN_FONT_SIZE)

    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.clear()

    content = slide_data.get("content", [])
    for i, item in enumerate(content):
        p = tf.paragraphs[0] if i == 0 and len(tf.paragraphs) == 1 else tf.add_paragraph()
        p.level = 0
        parts = re.split(r"(\*\*.*?\*\*)", str(item))
        for part in parts:
            if not part:
                continue
            run = p.add_run()
            if part.startswith("**") and part.endswith("**"):
                run.text = part[2:-2]
                run.font.bold = True
                run.font.color.rgb = RGBColor(0, 112, 192)
            else:
                run.text = part
            run.font.size = Pt(font_size_pt)
        spacing = max(3, font_size_pt * 0.3)
        p.space_before = Pt(spacing)
        p.space_after = Pt(spacing)
