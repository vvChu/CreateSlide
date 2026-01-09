from slide_engine import create_pptx
from pptx import Presentation
from pptx.util import Cm

def verify_layout_geometry():
    print("Testing Advanced Layout Config...")
    
    mock_data = {
        "title": "mixed case title",
        "slides": [
            {
                "title": "Slide 1",
                "content": ["Content 1"]
            }
        ]
    }
    
    pptx_io = create_pptx(mock_data)
    prs = Presentation(pptx_io)
    
    # 1. Verify Uppercase Title
    slide0 = prs.slides[0]
    title0 = slide0.shapes.title.text
    print(f"Slide 0 Title: '{title0}'")
    if title0 == "MIXED CASE TITLE":
        print("PASS: Title is uppercase.")
    else:
        print(f"FAIL: Title is '{title0}'")

    # Check Title Geometry on Slide 1 (Content Slide)
    # Note: slide 0 is Title Slide, usually different layout. Slide 1 is Content.
    # Check Title Geometry on Slide 1 (Content Slide)
    # Note: slide 0 is Title Slide, usually different layout. Slide 1 is Content.
    slide1 = prs.slides[1]
    title1 = slide1.shapes.title
    if title1:
        top_title_cm = title1.top.cm
        height_title_cm = title1.height.cm
        print(f"Title Top: {top_title_cm:.2f} cm (Expected ~0.5)")
        print(f"Title Height: {height_title_cm:.2f} cm (Expected ~3.5)")
        
        if 0.4 < top_title_cm < 0.6 and 3.4 < height_title_cm < 3.6:
             print("PASS: Title Geometry correct.")
        else:
             print("FAIL: Title Geometry mismatch.")
        
    # 2. Verify Body Geometry
    slide1 = prs.slides[1]
    body1 = None
    for shape in slide1.placeholders:
         if shape.placeholder_format.idx == 1:
             body1 = shape
             
    if body1:
        # Expected:
        # Left = 1.5 cm
        # Top = 5.5 cm
        
        left_cm = body1.left.cm
        top_cm = body1.top.cm
        
        print(f"Body Left: {left_cm:.2f} cm (Expected ~1.5)")
        if 1.4 < left_cm < 1.6 and 5.4 < top_cm < 5.6:
            print(f"PASS: Body Top matches precise layout (5.5 cm).")
        else:
            print(f"FAIL: Geometry mismatch. Top={top_cm}, Left={left_cm}")
    else:
        print("FAIL: Body placeholder not found.")

if __name__ == "__main__":
    verify_layout_geometry()
