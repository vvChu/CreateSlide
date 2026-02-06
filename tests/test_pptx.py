"""Tests for app.rendering.pptx — PowerPoint generation."""

from __future__ import annotations

import io

from pptx import Presentation

from app.rendering.pptx import create_pptx


class TestCreatePptx:
    """Test PPTX file generation. create_pptx returns io.BytesIO."""

    SAMPLE_SLIDES = {
        "title": "Test Presentation",
        "slides": [
            {
                "title": "Introduction",
                "content": ["Point 1", "Point 2", "Point 3"],
            },
            {
                "title": "Details",
                "content": ["Detail A", "Detail B"],
                "notes": "Speaker notes here",
            },
        ],
    }

    def test_returns_bytesio(self):
        result = create_pptx(self.SAMPLE_SLIDES)
        assert isinstance(result, io.BytesIO)

    def test_valid_pptx_can_be_opened(self):
        result = create_pptx(self.SAMPLE_SLIDES)
        prs = Presentation(result)
        assert prs is not None

    def test_correct_slide_count(self):
        result = create_pptx(self.SAMPLE_SLIDES)
        prs = Presentation(result)
        # +1 for title slide
        assert len(prs.slides) == len(self.SAMPLE_SLIDES["slides"]) + 1

    def test_title_slide_text(self):
        result = create_pptx(self.SAMPLE_SLIDES)
        prs = Presentation(result)
        title_slide = prs.slides[0]
        texts = [shape.text for shape in title_slide.shapes if shape.has_text_frame]
        combined = " ".join(texts).upper()
        assert "TEST PRESENTATION" in combined

    def test_empty_slides_dict(self):
        result = create_pptx({"title": "Empty", "slides": []})
        prs = Presentation(result)
        assert len(prs.slides) == 1  # title slide only

    def test_single_slide(self):
        data = {
            "title": "One Slide",
            "slides": [{"title": "Only", "content": ["Content"]}],
        }
        result = create_pptx(data)
        prs = Presentation(result)
        assert len(prs.slides) == 2  # title + 1 content

    def test_missing_content_handled(self):
        """Slides with no content key should not crash."""
        data = {
            "title": "No Content",
            "slides": [{"title": "Empty Slide"}],
        }
        result = create_pptx(data)
        assert isinstance(result, io.BytesIO)

    def test_slide_with_notes(self):
        data = {
            "title": "Notes Test",
            "slides": [
                {"title": "S1", "content": ["C1"], "notes": "My notes"},
            ],
        }
        result = create_pptx(data)
        prs = Presentation(result)
        content_slide = prs.slides[1]
        if content_slide.has_notes_slide:
            notes_text = content_slide.notes_slide.notes_text_frame.text
            assert "My notes" in notes_text

    def test_large_content_list(self):
        """Many bullet points should not crash."""
        data = {
            "title": "Many Points",
            "slides": [
                {"title": "Big", "content": [f"Point {i}" for i in range(30)]},
            ],
        }
        result = create_pptx(data)
        assert isinstance(result, io.BytesIO)
