from slide_engine import create_pptx
from pptx import Presentation
from pptx.dml.color import RGBColor

def verify_highlighting():
    print("Testing Content Highlighting...")
    
    mock_data = {
        "title": "Highlight Test",
        "slides": [
            {
                "title": "Slide 1",
                "content": ["Normal text **Bold Blue** Normal text."]
            }
        ]
    }
    
    pptx_io = create_pptx(mock_data)
    prs = Presentation(pptx_io)
    
    slide = prs.slides[1] # Content slide
    body = None
    for shape in slide.placeholders:
         if shape.placeholder_format.idx == 1:
             body = shape
             
    if body and body.has_text_frame:
        p = body.text_frame.paragraphs[0]
        print(f"Total Runs: {len(p.runs)}")
        
        for i, run in enumerate(p.runs):
            font = run.font
            is_bold = font.bold
            
            # Check color safely
            color_hex = "None"
            if font.color and font.color.type:
                 # It might be an RGB color or Theme color
                 try:
                     color_hex = str(font.color.rgb)
                 except:
                     color_hex = "Theme/Other"
                     
            print(f"Run {i}: Text='{run.text}', Bold={is_bold}, Color={color_hex}")
            
        # Assertions
        if len(p.runs) >= 3:
            r1 = p.runs[1]
            if r1.text == "Bold Blue" and r1.font.bold:
                print("PASS: Highlighting applied correctly.")
            else:
                print("FAIL: Middle run is not highlighted.")
        else:
            print("FAIL: Text was not split into runs.")

if __name__ == "__main__":
    verify_highlighting()
