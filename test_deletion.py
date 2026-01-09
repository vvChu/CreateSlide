from pptx import Presentation
import io

def test_deletion_logic():
    # 1. Create a dummy "Template" with 3 slides
    prs = Presentation()
    for i in range(3):
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"Existing Slide {i+1}"

    print(f"Original Slide Count: {len(prs.slides)}")

    # 2. Simulate "Naive" Deletion (What I did before)
    # The hypothesis: this deletes the slide ID but leaves the relationship, corrupting file.
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    
    # Try the method I used
    while len(xml_slides) > 0:
         xml_slides.remove(xml_slides[0])

    print(f"Post-Deletion Slide Count: {len(prs.slides)}")
    
    # 3. Save and "Re-Open" verification
    out = io.BytesIO()
    try:
        prs.save(out)
        print("Save Successful")
    except Exception as e:
        print(f"Save Failed: {e}")
        return

    # Verify if we can open it
    out.seek(0)
    try:
        new_prs = Presentation(out)
        print(f"Re-opened Slide Count: {len(new_prs.slides)}")
        if len(new_prs.slides) != 0:
            print("FAILURE: Slides persist or reappeared!")
    except Exception as e:
        print(f"Re-open Failed (Corruption confirmed?): {e}")

if __name__ == "__main__":
    test_deletion_logic()
