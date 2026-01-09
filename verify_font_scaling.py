from slide_engine import create_pptx
from pptx import Presentation
import io

def verify_font_scaling():
    print("Testing Font Scaling Logic...")
    
    # Mock Data with varying lengths
    mock_data = {
        "title": "Font Test",
        "slides": [
            {
                "title": "Slide 1: Short Content",
                "content": ["This is a short bullet."] * 2 # < 200 chars
            },
            {
                "title": "Slide 2: Medium Content",
                "content": ["This is a medium length bullet point repeated."] * 8 # ~300-400 chars
            },
            {
                "title": "Slide 3: Long Content",
                "content": ["This is a very long bullet point to trigger the smallest font size scaling logic."] * 15 # > 600 chars
            }
        ]
    }
    
    # specific template not needed for this logic check
    
    try:
        pptx_io = create_pptx(mock_data)
        prs = Presentation(pptx_io)
        
        # Check Slide 1 (Short)
        slide1 = prs.slides[1] # Index 1 is first content slide
        # Find body (heuristic: placeholder 1 or non-title)
        body1 = None
        for shape in slide1.placeholders:
             if shape.placeholder_format.idx == 1:
                 body1 = shape
        
        if body1 and body1.has_text_frame:
            print(f"Slide 1 (Short): Total Paragraphs = {len(body1.text_frame.paragraphs)}")
            for i, p in enumerate(body1.text_frame.paragraphs):
                 print(f"  Para {i}: Text='{p.text}', Runs={len(p.runs)}")
                 if p.runs:
                     print(f"    Run 0 Size: {p.runs[0].font.size.pt}pt (Expected 28.0)")
                 elif p.text:
                     # If text exists but no runs exposed? Rare case.
                     print(f"    Text exists but no runs found in this API view.")

        # Check Slide 3 (Long)
        slide3 = prs.slides[3] # Index 3 is 3rd content slide
        body3 = None
        for shape in slide3.placeholders:
             if shape.placeholder_format.idx == 1:
                 body3 = shape
        
        if body3 and body3.has_text_frame:
            print(f"Slide 3 (Long): Total Paragraphs = {len(body3.text_frame.paragraphs)}")
            for i, p in enumerate(body3.text_frame.paragraphs):
                 if len(p.runs) > 0:
                     font_pt = p.runs[0].font.size.pt
                     print(f"  Para {i} Run 0 Size: {font_pt}pt (Expected 16.0)")
                 else:
                     print(f"  Para {i}: No runs.")
            
        print("Verification Complete: No Crashes.")
        
    except Exception as e:
        print(f"Verification Failed: {e}")
        import traceback
        traceback.print_exc()

if __name__ == "__main__":
    verify_font_scaling()
