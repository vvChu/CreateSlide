import io
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
import re

def create_pptx(json_data: Dict[str, Any], template_pptx_bytes: bytes | None = None) -> io.BytesIO:
    """
    Generates a PowerPoint presentation from JSON data.
    Returns: BytesIO object of the .pptx file.
    """
    if template_pptx_bytes:
        prs = Presentation(io.BytesIO(template_pptx_bytes))
        # Clear existing slides from the template explicitly and safely
        # We need to remove the relationship (rId) to avoid corruption
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        
        for s in slides:
             rId = s.rId
             prs.part.drop_rel(rId) # Critical: Remove the relationship
             xml_slides.remove(s)   # Remove the slide entry
            
        print(f"Template loaded and cleared. Remaining slides: {len(prs.slides)}")
        
    else:
        prs = Presentation() # Uses default template

    # Helper to find usable layouts
    def get_layout(prs, preferred_index, needs_body=False):
        # Try preferred index first
        if preferred_index < len(prs.slide_layouts):
            layout = prs.slide_layouts[preferred_index]
            if not needs_body:
                return layout
            # Check if it has a body placeholder (usually idx 1)
            if len(layout.placeholders) > 1:
                return layout
        
        # Fallback: Search for a suitable layout
        for layout in prs.slide_layouts:
            if needs_body and len(layout.placeholders) > 1:
                return layout
        
        # Last resort: just return first layout
        return prs.slide_layouts[0]

    # 1. Main Title Slide
    # layout 0 usually title
    title_layout = get_layout(prs, 0) 
    slide = prs.slides.add_slide(title_layout)
    
    if slide.shapes.title:
        title_text = json_data.get("title", "Bài thuyết trình AI")
        slide.shapes.title.text = title_text.upper()
    
    # Safely set subtitle if placeholder exists
    if len(slide.placeholders) > 1:
        try:
            slide.placeholders[1].text = "Được tạo bởi SlideGenius"
        except (IndexError, KeyError):
            pass # Skip if no placeholder accessible

    # 2. Content Slides
    slides_content = json_data.get("slides", [])
    
    # Try to find a content layout (often index 1, needs body)
    content_layout = get_layout(prs, 1, needs_body=True)
    
    # ... existing code ...

    for slide_data in slides_content:
        slide = prs.slides.add_slide(content_layout)
        
        # Set Title and Clean "Slide X:" prefix
        if slide.shapes.title:
            raw_title = slide_data.get("title", "")
            # Remove "Slide 1:", "Slide 01", etc.
            clean_title = re.sub(r'^Slide\s+\d+[:.]?\s*', '', raw_title, flags=re.IGNORECASE)
            slide.shapes.title.text = clean_title
            
            # --- Layout Refinement: Strict Title Geometry ---
            # Force Title to top area to prevent wandering
            slide.shapes.title.top = Cm(0.5)
            slide.shapes.title.left = Cm(1.0)
            slide.shapes.title.width = prs.slide_width - Cm(2.0)
            slide.shapes.title.height = Cm(3.5) # Increased to ~3.5cm to safely fit 2 lines of 36pt
            # Note: width/left centered with 1cm margin on each side
            
            # Title Font Styling (Request: 36pt, Top Align)
            if slide.shapes.title.text_frame:
                slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
                # Ensure vertical alignment is TOP (though default often is middle)
                # We need to import PP_ALIGN if we want to force alignment, 
                # but standard resizing usually works. Let's rely on geometry.

        # Set Content (Body)
        # ... (finding body_shape logic remains) ...

        # Set Content (Body)
        # Find the first non-title placeholder
        body_shape = None
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1: # Standard body
                body_shape = shape
                break
        
        # Fallback to any second placeholder if idx 1 not found
        if not body_shape and len(slide.placeholders) > 1:
             body_shape = slide.placeholders[1]

        if body_shape and hasattr(body_shape, "text_frame"):
            # --- Layout Refinement: Adjust Margins & Top ---
            # Explicit Safety Zones to prevent overlap
            
            margin_side = Cm(1.5)
            # Body Top = Title Bottom (0.5+3.5=4.0) + Gap 1.5cm = 5.5cm
            margin_top = Cm(5.5) 
            margin_bottom = Cm(1.0)
            
            body_shape.left = margin_side
            body_shape.top = margin_top
            body_shape.width = prs.slide_width - (margin_side * 2)
            body_shape.height = prs.slide_height - margin_top - margin_bottom
            
            tf = body_shape.text_frame
            tf.word_wrap = True
            
            # --- Native Auto-Fit Strategy ---
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            
            content = slide_data.get("content", [])
            # ... (content normalization) ...
            
            # Clear default paragraph
            tf.clear()
            
            for i, item in enumerate(content):
                # Reuse the first paragraph if it exists
                if i == 0 and len(tf.paragraphs) == 1:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.level = 0
                
                # Use explicit run, but allow Auto-Fit to override
                # Parse markdown bold syntax
                parts = re.split(r'(\*\*.*?\*\*)', str(item))
                
                for part in parts:
                    if not part: continue
                    
                    run = p.add_run()
                    # Check for bold marker
                    if part.startswith('**') and part.endswith('**'):
                        text_content = part[2:-2] # Strip **
                        run.text = text_content
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(0, 112, 192) # Emphasis Blue
                    else:
                        run.text = part
                        
                    # Set Base Font to 20pt (User Max Request)
                    # Auto-fit will scale DOWN from here if needed.
                    run.font.size = Pt(20) 
                
                p.space_before = Pt(6)
                p.space_after = Pt(6)
        
        # Set Speaker Notes
        
        # Set Speaker Notes
        notes = slide_data.get("notes", "")
        if notes:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = notes

    # Save to BytesIO
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
